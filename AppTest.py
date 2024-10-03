from kivy.config import Config
Config.set('input', 'mouse', 'mouse,multitouch_on_demand')

from kivy.app import App
from kivy.uix.floatlayout import FloatLayout
from kivy.uix.boxlayout import BoxLayout
from kivy.uix.button import Button
from kivy.uix.label import Label
from kivy.uix.image import Image
from kivy.uix.filechooser import FileChooserListView
from kivy.uix.popup import Popup
from kivy.uix.scrollview import ScrollView
from kivy.clock import Clock
from kivy.graphics.texture import Texture
from kivy.core.window import Window
import cv2
import numpy as np
import pandas as pd
from keras.models import load_model # type: ignore
import PyPDF2
from pptx import Presentation
from PIL import Image as PILImage
import io
import os
import fitz  # PyMuPDF

class GestureRecognitionApp(App):
    def __init__(self, **kwargs):
        super().__init__(**kwargs)
        self.model = load_model(r"C:\Users\USER\Desktop\E2\HandGestureRecognitionSystemProject\model\resnetmodel.hdf5", compile=False)
        self.labels = pd.read_csv(r"C:\Users\USER\Documents\Dataset\jester-v1-labels.csv", header=None)
        self.capture = None
        self.buffer = []
        self.i = 1
        self.final_label = "No gesture"
        self.current_page = 0
        self.document = None
        self.zoom_level = 1.0

    def build(self):
        self.main_layout = FloatLayout()
        
        background = Image(source=r'C:\Users\USER\Desktop\E2\HandGestureRecognitionSystemProject\bgr4.jpeg', allow_stretch=True, keep_ratio=False)
        self.main_layout.add_widget(background)
        
        content_layout = BoxLayout(orientation='vertical', padding=10, spacing=10)
        
        top_layout = BoxLayout(size_hint_y=0.1)
        self.title_label = Label(text='Gesture Recognition System', font_size=24)
        self.status_label = Label(text='Status: Ready', font_size=18)
        top_layout.add_widget(self.title_label)
        top_layout.add_widget(self.status_label)
        
        middle_layout = BoxLayout(size_hint_y=0.6)
        self.camera_feed = Image(size_hint_x=0.5)
        self.document_view = ScrollView(size_hint_x=0.5)
        self.document_image = Image(allow_stretch=True, keep_ratio=True)
        self.document_view.add_widget(self.document_image)
        middle_layout.add_widget(self.camera_feed)
        middle_layout.add_widget(self.document_view)
        
        self.recognized_gesture = Label(text='Recognized Gesture: None', font_size=20, size_hint_y=0.1)
        
        bottom_layout = BoxLayout(size_hint_y=0.2, spacing=10)
        start_button = Button(text='Start Recognition', on_press=self.start_recognition)
        stop_button = Button(text='Stop Recognition', on_press=self.stop_recognition)
        browse_button = Button(text='Browse Documents', on_press=self.show_file_chooser)
        bottom_layout.add_widget(start_button)
        bottom_layout.add_widget(stop_button)
        bottom_layout.add_widget(browse_button)
        
        content_layout.add_widget(top_layout)
        content_layout.add_widget(middle_layout)
        content_layout.add_widget(self.recognized_gesture)
        content_layout.add_widget(bottom_layout)
        
        self.main_layout.add_widget(content_layout)

        Window.bind(on_dropfile=self.on_file_drop)
        
        return self.main_layout

    def show_file_chooser(self, instance):
        content = BoxLayout(orientation='vertical')
        file_chooser = FileChooserListView(filters=['*.pdf', '*.pptx'])
        content.add_widget(file_chooser)
        
        select_button = Button(text='Select', size_hint_y=None, height=50)
        select_button.bind(on_release=lambda x: self.load_document(file_chooser.selection) or self.dismiss_popup())
        content.add_widget(select_button)
        
        self._popup = Popup(title="Select a document", content=content, size_hint=(0.9, 0.9))
        self._popup.open()

    def dismiss_popup(self):
        self._popup.dismiss()

    def load_document(self, selection):
        if selection:
            try:
                file_path = selection[0]
                if file_path.lower().endswith('.pdf'):
                    self.document = fitz.open(file_path)
                    self.render_pdf_page(0)
                elif file_path.lower().endswith('.pptx'):
                    self.document = Presentation(file_path)
                    self.render_pptx_slide(0)
                self.status_label.text = f'Status: Loaded {os.path.basename(file_path)}'
                self.current_page = 0
            except Exception as e:
                self.status_label.text = f'Error: Failed to load document. {str(e)}'

    def render_pdf_page(self, page_num):
        if self.document and page_num < len(self.document):
            page = self.document[page_num]
            pix = page.get_pixmap()
            img = PILImage.frombytes("RGB", [pix.width, pix.height], pix.samples)
            
            # Convert PIL image to Kivy texture
            data = io.BytesIO()
            img.save(data, format='png')
            data.seek(0)
            
            texture = Texture.create(size=(pix.width, pix.height))
            texture.blit_buffer(data.read(), colorfmt='rgb', bufferfmt='ubyte')
            
            self.document_image.texture = texture
            self.document_image.size = texture.size

    def render_pptx_slide(self, slide_num):
        if self.document and slide_num < len(self.document.slides):
            slide = self.document.slides[slide_num]
            
            # Convert slide to image (you may need to adjust this part)
            img_stream = io.BytesIO()
            slide.export(img_stream, 'PNG')
            img_stream.seek(0)
            
            # Convert to Kivy texture
            img = PILImage.open(img_stream)
            texture = Texture.create(size=img.size)
            texture.blit_buffer(img.tobytes(), colorfmt='rgb', bufferfmt='ubyte')
            
            self.document_image.texture = texture
            self.document_image.size = texture.size

    def start_recognition(self, instance):
        self.capture = cv2.VideoCapture(0)
        self.capture.set(cv2.CAP_PROP_FRAME_WIDTH, 400)
        self.capture.set(cv2.CAP_PROP_FRAME_HEIGHT, 400)
        Clock.schedule_interval(self.update, 1.0/30.0)
        self.status_label.text = 'Status: Running'

    def stop_recognition(self, instance):
        Clock.unschedule(self.update)
        if self.capture:
            self.capture.release()
        self.camera_feed.texture = None
        self.status_label.text = 'Status: Stopped'
        self.recognized_gesture.text = 'Recognized Gesture: None'

    def update(self, dt):
        ret, frame = self.capture.read()
        if ret:
            image = cv2.resize(frame, (96, 64))
            image = image / 255.0
            self.buffer.append(image)
            if self.i % 16 == 0:
                self.buffer = np.expand_dims(self.buffer, 0)
                predicted_value = np.argmax(self.model.predict(self.buffer))
                cls = self.labels.iloc[predicted_value]
                self.final_label = cls.iloc[0]
                self.perform_action(predicted_value)
                self.buffer = []

            self.i += 1
            
            buf1 = cv2.flip(frame, 0)
            buf = buf1.tostring()
            texture1 = Texture.create(size=(frame.shape[1], frame.shape[0]), colorfmt='bgr')
            texture1.blit_buffer(buf, colorfmt='bgr', bufferfmt='ubyte')
            self.camera_feed.texture = texture1
            self.recognized_gesture.text = f'Recognized Gesture: {self.final_label}'

    def perform_action(self, predicted_value):
        if self.document:
            if predicted_value == 0:  # Swiping Left
                self.next_page()
            elif predicted_value == 1:  # Swiping Right
                self.previous_page()
            elif predicted_value == 2:  # Swiping Down
                self.scroll_down()
            elif predicted_value == 3:  # Swiping Up
                self.scroll_up()
            elif predicted_value == 18:  # Zooming in with two fingers
                self.zoom_in()
            elif predicted_value == 19:  # Zooming out with two fingers
                self.zoom_out()

    def next_page(self):
        if isinstance(self.document, fitz.Document):
            if self.current_page < len(self.document) - 1:
                self.current_page += 1
                self.render_pdf_page(self.current_page)
        elif isinstance(self.document, Presentation):
            if self.current_page < len(self.document.slides) - 1:
                self.current_page += 1
                self.render_pptx_slide(self.current_page)

    def previous_page(self):
        if self.current_page > 0:
            self.current_page -= 1
            if isinstance(self.document, fitz.Document):
                self.render_pdf_page(self.current_page)
            elif isinstance(self.document, Presentation):
                self.render_pptx_slide(self.current_page)

    def scroll_down(self):
        self.document_view.scroll_y = max(self.document_view.scroll_y - 0.1, 0)

    def scroll_up(self):
        self.document_view.scroll_y = min(self.document_view.scroll_y + 0.1, 1)

    def zoom_in(self):
        self.zoom_level *= 1.1
        self.update_zoom()

    def zoom_out(self):
        self.zoom_level /= 1.1
        self.update_zoom()

    def update_zoom(self):
        self.document_image.size = (Window.width * self.zoom_level, Window.height * self.zoom_level)
        self.document_image.reload()

    def on_file_drop(self, window, file_path):
        file_path = file_path.decode("utf-8")
        if file_path.lower().endswith(('.pdf', '.pptx')):
            self.load_document([file_path])
        else:
            self.status_label.text = "Status: Invalid file type dropped."

if __name__ == '__main__':
    GestureRecognitionApp().run()